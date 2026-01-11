# -*- coding: utf-8 -*-
"""
VIBE CODING STARTER - Module 1 Presentation Generator
Brand: CYBER-ARCHITECTURE v3.0
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================================
# BRAND COLORS
# ============================================================================
CYBER_ACID = RgbColor(0x00, 0xFF, 0x88)      # #00FF88 - CTA, accents
CYBER_VOID = RgbColor(0x03, 0x03, 0x03)      # #030303 - Background
CYBER_SURFACE = RgbColor(0x0A, 0x0A, 0x0A)   # #0A0A0A - Cards
SIGNAL_RED = RgbColor(0xFF, 0x33, 0x66)      # #FF3366 - Warnings
HOLO_WHITE = RgbColor(0xFF, 0xFF, 0xFF)      # #FFFFFF - Main text
TECH_GRAY = RgbColor(0x9C, 0xA3, 0xAF)       # #9CA3AF - Secondary text

# Semi-transparent versions (simulated with solid colors for PPTX)
ACID_15 = RgbColor(0x0D, 0x2B, 0x1B)         # ~15% green on black
ACID_30 = RgbColor(0x15, 0x4D, 0x2F)         # ~30% green on black
RED_10 = RgbColor(0x2B, 0x0D, 0x15)          # ~10% red on black
GREEN_10 = RgbColor(0x0D, 0x2B, 0x1B)        # ~10% green on black
SURFACE_90 = RgbColor(0x0C, 0x0C, 0x0C)      # Surface with slight transparency

# ============================================================================
# DIMENSIONS (16:9 - 1920x1080 in EMU)
# ============================================================================
SLIDE_WIDTH = Inches(13.333)   # 1920px at 144dpi ~ 13.33"
SLIDE_HEIGHT = Inches(7.5)     # 1080px at 144dpi ~ 7.5"

# Margins
MARGIN = Inches(0.4)           # ~32px minimum margin
CONTENT_LEFT = Inches(0.5)
CONTENT_TOP = Inches(0.8)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_slide_background(slide, color):
    """Set solid background color for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text,
                 font_name='Arial', font_size=28, font_bold=False,
                 font_color=HOLO_WHITE, alignment=PP_ALIGN.LEFT,
                 vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = font_bold
    p.font.color.rgb = font_color
    p.alignment = alignment

    tf.vertical_anchor = vertical_anchor

    return txBox

def add_shape_with_text(slide, left, top, width, height, text,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                        fill_color=None, line_color=None, line_width=Pt(1),
                        font_name='Arial', font_size=28, font_bold=False,
                        font_color=HOLO_WHITE, alignment=PP_ALIGN.CENTER):
    """Add a shape with centered text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = font_bold
    p.font.color.rgb = font_color
    p.alignment = alignment
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return shape

def add_rectangle(slide, left, top, width, height, fill_color=None,
                  line_color=None, line_width=Pt(1)):
    """Add a simple rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape

def add_line(slide, start_x, start_y, end_x, end_y, color=CYBER_ACID, width=Pt(3)):
    """Add a line"""
    line = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    line.line.color.rgb = color
    line.line.width = width
    return line

def add_badge(slide, text, right_offset=Inches(0.5), top=Inches(0.4)):
    """Add a lesson badge in top-right corner"""
    badge_width = Inches(1.8)
    badge_height = Inches(0.45)
    left = SLIDE_WIDTH - badge_width - right_offset

    badge = add_shape_with_text(
        slide, left, top, badge_width, badge_height, text,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        fill_color=CYBER_ACID,
        font_name='Courier New', font_size=18, font_bold=True,
        font_color=CYBER_VOID
    )
    return badge

def add_footer(slide, text="VIBE CODING STARTER • МОДУЛЬ 1"):
    """Add footer text at bottom"""
    add_text_box(
        slide, CONTENT_LEFT, SLIDE_HEIGHT - Inches(0.6),
        Inches(6), Inches(0.4), text,
        font_name='Courier New', font_size=14,
        font_color=TECH_GRAY
    )

def add_card(slide, left, top, width, height, fill_color=SURFACE_90,
             border_color=None, border_width=Pt(1)):
    """Add a card (rounded rectangle with subtle styling)"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color

    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.color.rgb = RgbColor(0x20, 0x20, 0x20)
        card.line.width = Pt(1)

    return card

def add_number_indicator(slide, left, top, number, size=Inches(0.7)):
    """Add a square number indicator with green accent"""
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACID_15
    shape.line.color.rgb = ACID_30
    shape.line.width = Pt(1)

    # Number text
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.name = 'Arial Black'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYBER_ACID
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return shape

def add_step_circle(slide, left, top, number, size=Inches(0.5)):
    """Add a circular step number"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYBER_ACID
    circle.line.fill.background()

    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.name = 'Arial Black'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYBER_VOID
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return circle

def add_checkbox(slide, left, top, checked=True, size=Inches(0.35)):
    """Add a checkbox with checkmark"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size
    )
    box.fill.background()
    box.line.color.rgb = CYBER_ACID
    box.line.width = Pt(2)

    if checked:
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "+"
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CYBER_ACID
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return box

# ============================================================================
# SLIDE CREATORS
# ============================================================================

def create_slide_1_title(prs):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Subtitle: VIBE CODING STARTER
    add_text_box(
        slide, Inches(0), Inches(1.8),
        SLIDE_WIDTH, Inches(0.5), "VIBE CODING STARTER",
        font_name='Courier New', font_size=22,
        font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
    )

    # Main title: БЫСТРЫЙ СТАРТ
    add_text_box(
        slide, Inches(0), Inches(2.5),
        SLIDE_WIDTH, Inches(1.2), "БЫСТРЫЙ",
        font_name='Arial Black', font_size=96,
        font_color=HOLO_WHITE, alignment=PP_ALIGN.CENTER,
        font_bold=True
    )

    add_text_box(
        slide, Inches(0), Inches(3.5),
        SLIDE_WIDTH, Inches(1.2), "СТАРТ",
        font_name='Arial Black', font_size=96,
        font_color=HOLO_WHITE, alignment=PP_ALIGN.CENTER,
        font_bold=True
    )

    # Tagline
    add_text_box(
        slide, Inches(0), Inches(4.8),
        SLIDE_WIDTH, Inches(0.6), "От нуля до первого проекта за 4 часа",
        font_name='Arial', font_size=32,
        font_color=TECH_GRAY, alignment=PP_ALIGN.CENTER
    )

    # Accent line
    line_width = Inches(2.5)
    line_left = (SLIDE_WIDTH - line_width) / 2
    add_rectangle(
        slide, line_left, Inches(5.6),
        line_width, Pt(3),
        fill_color=CYBER_ACID
    )

    return slide

def create_slide_2_overview(prs):
    """Slide 2: What awaits you in the module"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(10), Inches(0.9), "ЧТО ВАС ЖДЁТ В МОДУЛЕ",
        font_name='Arial Black', font_size=56,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - Stats cards
    stats = [
        ("6", "Уроков", "От основ до первого проекта"),
        ("~4ч", "Часа контента", "Концентрированные знания"),
        ("1", "Готовый проект", "Калькулятор в портфолио")
    ]

    card_top = Inches(1.6)
    for i, (num, title, desc) in enumerate(stats):
        # Card background
        card = add_card(
            slide, CONTENT_LEFT, card_top + Inches(i * 1.3),
            Inches(5.5), Inches(1.1)
        )

        # Number indicator
        add_number_indicator(
            slide, CONTENT_LEFT + Inches(0.2),
            card_top + Inches(i * 1.3) + Inches(0.2),
            num, size=Inches(0.7)
        )

        # Title
        add_text_box(
            slide, CONTENT_LEFT + Inches(1.1),
            card_top + Inches(i * 1.3) + Inches(0.15),
            Inches(4), Inches(0.45), title,
            font_name='Arial', font_size=26,
            font_color=HOLO_WHITE, font_bold=True
        )

        # Description
        add_text_box(
            slide, CONTENT_LEFT + Inches(1.1),
            card_top + Inches(i * 1.3) + Inches(0.55),
            Inches(4), Inches(0.4), desc,
            font_name='Arial', font_size=20,
            font_color=TECH_GRAY
        )

    # Right column - Result card
    result_card = add_card(
        slide, Inches(7), Inches(1.6),
        Inches(5.5), Inches(2.5),
        fill_color=ACID_15, border_color=ACID_30
    )

    add_text_box(
        slide, Inches(7.3), Inches(1.8),
        Inches(5), Inches(0.4), "РЕЗУЛЬТАТ МОДУЛЯ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.4),
        Inches(5), Inches(0.5), "Базовые навыки",
        font_name='Arial', font_size=28,
        font_color=CYBER_ACID, font_bold=True
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.9),
        Inches(5), Inches(0.5), "AI-программирования",
        font_name='Arial', font_size=28,
        font_color=CYBER_ACID, font_bold=True
    )

    add_text_box(
        slide, Inches(7.3), Inches(3.5),
        Inches(5), Inches(0.4), "и первый работающий проект",
        font_name='Arial', font_size=24,
        font_color=HOLO_WHITE
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_3_program(prs):
    """Slide 3: Module program"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(10), Inches(0.9), "ПРОГРАММА МОДУЛЯ",
        font_name='Arial Black', font_size=56,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Program items
    lessons = [
        ("1.1", "Революция в программировании", "до 20 мин", False),
        ("1.2", "Установка за 15 минут", "до 30 мин", False),
        ("1.3", "Git — твоя страховка", "до 40 мин", False),
        ("1.4", "Анатомия промпта", "до 20 мин", False),
        ("1.5", "Первый проект — Калькулятор", "до 60 мин", True),  # Highlighted
        ("1.6", "Что дальше + Бонус", "до 20 мин", False),
    ]

    row_height = Inches(0.75)
    start_top = Inches(1.5)

    for i, (num, title, time, highlighted) in enumerate(lessons):
        top = start_top + Inches(i * 0.85)

        # Row background
        if highlighted:
            row_bg = add_card(
                slide, CONTENT_LEFT, top,
                Inches(11), row_height,
                fill_color=ACID_15, border_color=ACID_30
            )
        else:
            row_bg = add_card(
                slide, CONTENT_LEFT, top,
                Inches(11), row_height
            )

        # Left accent line
        add_rectangle(
            slide, CONTENT_LEFT, top,
            Pt(4), row_height,
            fill_color=CYBER_ACID
        )

        # Lesson number
        add_text_box(
            slide, CONTENT_LEFT + Inches(0.2), top + Inches(0.2),
            Inches(0.8), Inches(0.4), num,
            font_name='Courier New', font_size=22,
            font_color=CYBER_ACID, font_bold=True
        )

        # Lesson title
        add_text_box(
            slide, CONTENT_LEFT + Inches(1.1), top + Inches(0.2),
            Inches(7), Inches(0.4), title,
            font_name='Arial', font_size=24,
            font_color=HOLO_WHITE
        )

        # Time
        time_color = CYBER_ACID if highlighted else TECH_GRAY
        add_text_box(
            slide, Inches(10), top + Inches(0.2),
            Inches(2), Inches(0.4), time,
            font_name='Courier New', font_size=20,
            font_color=time_color, alignment=PP_ALIGN.RIGHT
        )

    # Footer with total
    add_footer(slide)

    add_text_box(
        slide, Inches(9), SLIDE_HEIGHT - Inches(0.6),
        Inches(3.5), Inches(0.4), "ИТОГО: ~4 ЧАСА",
        font_name='Courier New', font_size=16,
        font_color=CYBER_ACID, alignment=PP_ALIGN.RIGHT
    )

    return slide

def create_slide_4_lesson_1_1(prs):
    """Slide 4: Lesson 1.1 - What is Vibe Coding"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.1")

    # Time indicator
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(3), Inches(0.4), "ДО 20 МИНУТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(10), Inches(0.8), "ЧТО ТАКОЕ VIBE CODING",
        font_name='Arial Black', font_size=52,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - What you'll learn
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(4), Inches(0.5), "Вы узнаете:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    checklist = [
        "Определение: AI-парное программирование",
        "Почему это работает сейчас",
        "Демонстрация на примере",
        "Отличие от классического"
    ]

    for i, item in enumerate(checklist):
        top = Inches(2.5) + Inches(i * 0.65)
        add_checkbox(slide, CONTENT_LEFT, top)
        add_text_box(
            slide, CONTENT_LEFT + Inches(0.5), top - Inches(0.05),
            Inches(5), Inches(0.5), item,
            font_name='Arial', font_size=22,
            font_color=HOLO_WHITE
        )

    # Right column - Formula card
    formula_card = add_card(
        slide, Inches(7), Inches(1.9),
        Inches(5.5), Inches(4.2),
        fill_color=SURFACE_90, border_color=ACID_30
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.4), "ГЛАВНАЯ ФОРМУЛА",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Formula elements
    formula_items = ["ИДЕЯ", "+", "КОНТЕКСТ", "+", "AI"]
    formula_top = Inches(2.7)

    for i, item in enumerate(formula_items):
        color = HOLO_WHITE if item != "+" else CYBER_ACID
        size = 28 if item != "+" else 24
        add_text_box(
            slide, Inches(8.5), formula_top + Inches(i * 0.45),
            Inches(3), Inches(0.4), item,
            font_name='Arial Black', font_size=size,
            font_color=color, alignment=PP_ALIGN.CENTER
        )

    # Divider line
    add_rectangle(
        slide, Inches(8), Inches(5),
        Inches(4), Pt(2),
        fill_color=HOLO_WHITE
    )

    # Result
    add_text_box(
        slide, Inches(8.5), Inches(5.3),
        Inches(3), Inches(0.5), "= ЧИСТЫЙ КОД",
        font_name='Arial Black', font_size=26,
        font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_5_lesson_1_2(prs):
    """Slide 5: Lesson 1.2 - Environment Setup"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.2")

    # Time indicator
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(3), Inches(0.4), "ДО 30 МИНУТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(10), Inches(0.8), "УСТАНОВКА ОКРУЖЕНИЯ",
        font_name='Arial Black', font_size=52,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - Steps
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(4), Inches(0.5), "Что установим:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    steps = [
        ("1", "VS Code", "Редактор кода"),
        ("2", "Kilo Code", "AI-расширение"),
        ("3", "Первый запуск", "Проверка работы AI")
    ]

    for i, (num, title, desc) in enumerate(steps):
        top = Inches(2.5) + Inches(i * 1.1)
        add_step_circle(slide, CONTENT_LEFT, top)

        add_text_box(
            slide, CONTENT_LEFT + Inches(0.7), top - Inches(0.05),
            Inches(4), Inches(0.4), title,
            font_name='Arial', font_size=24,
            font_color=HOLO_WHITE, font_bold=True
        )

        add_text_box(
            slide, CONTENT_LEFT + Inches(0.7), top + Inches(0.35),
            Inches(4), Inches(0.35), desc,
            font_name='Arial', font_size=20,
            font_color=TECH_GRAY
        )

    # Right column - Result card
    result_card = add_card(
        slide, Inches(7), Inches(1.9),
        Inches(5.5), Inches(3.5),
        fill_color=SURFACE_90, border_color=ACID_30
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.4), "РЕЗУЛЬТАТ УРОКА",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.7),
        Inches(5), Inches(0.8), "Полностью настроенное\nрабочее окружение",
        font_name='Arial', font_size=26,
        font_color=HOLO_WHITE
    )

    # Free badge
    free_badge = add_shape_with_text(
        slide, Inches(7.5), Inches(3.9),
        Inches(3), Inches(0.5), "100% БЕСПЛАТНО",
        fill_color=CYBER_ACID,
        font_name='Courier New', font_size=18, font_bold=True,
        font_color=CYBER_VOID
    )

    add_text_box(
        slide, Inches(7.3), Inches(4.7),
        Inches(5), Inches(0.4), "Все инструменты — open source",
        font_name='Arial', font_size=18,
        font_color=TECH_GRAY
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_6_lesson_1_3(prs):
    """Slide 6: Lesson 1.3 - Git"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.3")

    # Time indicator
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(3), Inches(0.4), "ДО 40 МИНУТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(10), Inches(0.8), "GIT — ТВОЯ СТРАХОВКА",
        font_name='Arial Black', font_size=52,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(4), Inches(0.5), "Зачем это нужно:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    # Without Git card (red)
    without_card = add_card(
        slide, CONTENT_LEFT, Inches(2.5),
        Inches(5.5), Inches(0.9),
        fill_color=RED_10, border_color=SIGNAL_RED
    )

    add_text_box(
        slide, CONTENT_LEFT + Inches(0.2), Inches(2.55),
        Inches(5), Inches(0.35), "БЕЗ GIT",
        font_name='Courier New', font_size=16,
        font_color=SIGNAL_RED
    )

    add_text_box(
        slide, CONTENT_LEFT + Inches(0.2), Inches(2.95),
        Inches(5), Inches(0.35), "Одна ошибка — код потерян",
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    # With Git card (green)
    with_card = add_card(
        slide, CONTENT_LEFT, Inches(3.6),
        Inches(5.5), Inches(0.9),
        fill_color=GREEN_10, border_color=CYBER_ACID
    )

    add_text_box(
        slide, CONTENT_LEFT + Inches(0.2), Inches(3.65),
        Inches(5), Inches(0.35), "С GIT",
        font_name='Courier New', font_size=16,
        font_color=CYBER_ACID
    )

    add_text_box(
        slide, CONTENT_LEFT + Inches(0.2), Inches(4.05),
        Inches(5), Inches(0.35), "Всегда можно откатиться",
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    # Right column - Commands card
    cmd_card = add_card(
        slide, Inches(7), Inches(1.9),
        Inches(5.5), Inches(3.8),
        fill_color=SURFACE_90, border_color=ACID_30
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.4), "4 КОМАНДЫ НА СТАРТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    commands = [
        ("git init", "Создать"),
        ("git add .", "Добавить"),
        ("git commit", "Сохранить"),
        ("git push", "В облако")
    ]

    for i, (cmd, desc) in enumerate(commands):
        top = Inches(2.7) + Inches(i * 0.7)

        # Command badge
        cmd_bg = add_shape_with_text(
            slide, Inches(7.3), top,
            Inches(2.2), Inches(0.45), cmd,
            fill_color=ACID_15,
            font_name='Courier New', font_size=18,
            font_color=CYBER_ACID
        )

        # Description
        add_text_box(
            slide, Inches(9.7), top + Inches(0.05),
            Inches(2.5), Inches(0.4), desc,
            font_name='Arial', font_size=20,
            font_color=HOLO_WHITE
        )

    # Result line
    add_text_box(
        slide, CONTENT_LEFT, Inches(5.2),
        Inches(11), Inches(0.4),
        "РЕЗУЛЬТАТ: Умение сохранять и откатывать изменения",
        font_name='Arial', font_size=20,
        font_color=CYBER_ACID
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_7_lesson_1_4(prs):
    """Slide 7: Lesson 1.4 - Prompt Anatomy"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.4")

    # Time indicator
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(3), Inches(0.4), "ДО 20 МИНУТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(10), Inches(0.8), "АНАТОМИЯ ПРОМПТА",
        font_name='Arial Black', font_size=52,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - Structure
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(4), Inches(0.5), "Структура промпта:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    structure = [
        ("КОНТЕКСТ", "Что за проект"),
        ("ЗАДАЧА", "Что конкретно сделать"),
        ("ОГРАНИЧЕНИЯ", "Что НЕ делать"),
        ("РЕЗУЛЬТАТ", "Как должно выглядеть")
    ]

    for i, (title, desc) in enumerate(structure):
        top = Inches(2.5) + Inches(i * 0.85)

        # Left accent
        add_rectangle(
            slide, CONTENT_LEFT, top,
            Pt(4), Inches(0.7),
            fill_color=CYBER_ACID
        )

        # Title
        add_text_box(
            slide, CONTENT_LEFT + Inches(0.15), top + Inches(0.05),
            Inches(4), Inches(0.35), title,
            font_name='Courier New', font_size=18,
            font_color=CYBER_ACID
        )

        # Description
        add_text_box(
            slide, CONTENT_LEFT + Inches(0.15), top + Inches(0.38),
            Inches(4), Inches(0.3), desc,
            font_name='Arial', font_size=18,
            font_color=TECH_GRAY
        )

    # Right column - Examples
    # Bad example
    bad_card = add_card(
        slide, Inches(7), Inches(2),
        Inches(5.5), Inches(1.3),
        fill_color=RED_10, border_color=SIGNAL_RED
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.35), "ПЛОХО",
        font_name='Courier New', font_size=16,
        font_color=SIGNAL_RED
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.55),
        Inches(5), Inches(0.5), '"Сделай мне сайт"',
        font_name='Arial', font_size=22,
        font_color=HOLO_WHITE
    )

    # Good example
    good_card = add_card(
        slide, Inches(7), Inches(3.6),
        Inches(5.5), Inches(2),
        fill_color=GREEN_10, border_color=CYBER_ACID
    )

    add_text_box(
        slide, Inches(7.3), Inches(3.7),
        Inches(5), Inches(0.35), "ХОРОШО",
        font_name='Courier New', font_size=16,
        font_color=CYBER_ACID
    )

    add_text_box(
        slide, Inches(7.3), Inches(4.15),
        Inches(5), Inches(1.2),
        '"Создай HTML калькулятор.\nИспользуй Vanilla JS.\nБез библиотек."',
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    # Result line
    add_text_box(
        slide, CONTENT_LEFT, Inches(6),
        Inches(11), Inches(0.4),
        "РЕЗУЛЬТАТ: Умение формулировать задачи для AI",
        font_name='Arial', font_size=20,
        font_color=CYBER_ACID
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_8_lesson_1_5(prs):
    """Slide 8: Lesson 1.5 - First Project"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.5")

    # Time indicator with highlight
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(5), Inches(0.4), "ДО 60 МИНУТ • ГЛАВНЫЙ УРОК",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(11), Inches(0.8), "ПЕРВЫЙ ПРОЕКТ — КАЛЬКУЛЯТОР",
        font_name='Arial Black', font_size=48,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - Steps
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(5), Inches(0.5), "Полный цикл разработки:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    steps = [
        ("1", "Идея", "определяем что делаем"),
        ("2", "Промпт", "формулируем задачу"),
        ("3", "Код", "AI генерирует"),
        ("4", "Тест", "проверяем")
    ]

    for i, (num, title, desc) in enumerate(steps):
        top = Inches(2.5) + Inches(i * 0.9)

        add_step_circle(slide, CONTENT_LEFT, top)

        add_text_box(
            slide, CONTENT_LEFT + Inches(0.7), top - Inches(0.05),
            Inches(2), Inches(0.4), title,
            font_name='Arial', font_size=24,
            font_color=HOLO_WHITE, font_bold=True
        )

        add_text_box(
            slide, CONTENT_LEFT + Inches(0.7), top + Inches(0.35),
            Inches(4), Inches(0.35), desc,
            font_name='Arial', font_size=18,
            font_color=TECH_GRAY
        )

        # Arrow down (except last)
        if i < 3:
            add_text_box(
                slide, CONTENT_LEFT + Inches(0.15), top + Inches(0.65),
                Inches(0.3), Inches(0.3), "|",
                font_name='Arial', font_size=20,
                font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
            )

    # Right column - Result card (highlighted)
    result_card = add_card(
        slide, Inches(7), Inches(1.9),
        Inches(5.5), Inches(3.8),
        fill_color=ACID_15, border_color=CYBER_ACID
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.4), "РЕЗУЛЬТАТ УРОКА",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Calculator icon placeholder
    add_text_box(
        slide, Inches(8.5), Inches(2.8),
        Inches(3), Inches(0.8), "[CALC]",
        font_name='Courier New', font_size=36,
        font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(7.3), Inches(3.8),
        Inches(5), Inches(0.5), "Работающий калькулятор",
        font_name='Arial', font_size=26,
        font_color=HOLO_WHITE, font_bold=True,
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(7.3), Inches(4.5),
        Inches(5), Inches(0.5), "Первый проект в портфолио!",
        font_name='Arial', font_size=22,
        font_color=CYBER_ACID,
        alignment=PP_ALIGN.CENTER
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_9_lesson_1_6(prs):
    """Slide 9: Lesson 1.6 - What's Next"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Badge
    add_badge(slide, "УРОК 1.6")

    # Time indicator
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(3), Inches(0.4), "ДО 20 МИНУТ",
        font_name='Courier New', font_size=18,
        font_color=CYBER_ACID
    )

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.9),
        Inches(10), Inches(0.8), "ЧТО ДАЛЬШЕ + БОНУС",
        font_name='Arial Black', font_size=52,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Left column - Checklist
    add_text_box(
        slide, CONTENT_LEFT, Inches(1.9),
        Inches(5), Inches(0.5), "Чеклист самопроверки:",
        font_name='Arial', font_size=24,
        font_color=TECH_GRAY
    )

    checklist = [
        "VS Code + Kilo Code установлены",
        "Git настроен и работает",
        "Знаю структуру промпта",
        "Калькулятор работает"
    ]

    for i, item in enumerate(checklist):
        top = Inches(2.5) + Inches(i * 0.65)
        add_checkbox(slide, CONTENT_LEFT, top)
        add_text_box(
            slide, CONTENT_LEFT + Inches(0.5), top - Inches(0.05),
            Inches(5), Inches(0.5), item,
            font_name='Arial', font_size=20,
            font_color=HOLO_WHITE
        )

    # Next step card
    next_card = add_card(
        slide, CONTENT_LEFT, Inches(5.2),
        Inches(5.5), Inches(0.8),
        fill_color=SURFACE_90, border_color=ACID_30
    )

    add_text_box(
        slide, CONTENT_LEFT + Inches(0.2), Inches(5.35),
        Inches(5), Inches(0.5), "Следующий шаг: Модуль 2: BUILDER ->",
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    # Right column - Bonus card
    bonus_card = add_card(
        slide, Inches(7), Inches(1.9),
        Inches(5.5), Inches(3.8),
        fill_color=RED_10, border_color=SIGNAL_RED
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.1),
        Inches(5), Inches(0.4), "БОНУС К МОДУЛЮ",
        font_name='Courier New', font_size=18,
        font_color=SIGNAL_RED
    )

    add_text_box(
        slide, Inches(7.3), Inches(2.7),
        Inches(5), Inches(0.6), "10 промптов",
        font_name='Arial', font_size=32,
        font_color=HOLO_WHITE, font_bold=True
    )

    add_text_box(
        slide, Inches(7.3), Inches(3.3),
        Inches(5), Inches(0.4), "Готовые шаблоны",
        font_name='Arial', font_size=22,
        font_color=TECH_GRAY
    )

    add_text_box(
        slide, Inches(7.3), Inches(4),
        Inches(5), Inches(0.4), "+ Шаблон Memory Bank",
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    add_text_box(
        slide, Inches(7.3), Inches(4.5),
        Inches(5), Inches(0.4), "+ Чеклист установки",
        font_name='Arial', font_size=20,
        font_color=HOLO_WHITE
    )

    # Footer
    add_footer(slide)

    return slide

def create_slide_10_results(prs):
    """Slide 10: Module Results"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Title
    add_text_box(
        slide, CONTENT_LEFT, Inches(0.5),
        Inches(10), Inches(0.9), "РЕЗУЛЬТАТ МОДУЛЯ",
        font_name='Arial Black', font_size=56,
        font_color=HOLO_WHITE, font_bold=True
    )

    # Results grid (2x2)
    results = [
        ("[ENV]", "Настроенное окружение", "VS Code + Kilo", False),
        ("[GIT]", "Базовые навыки Git", "Сохранение/откат", False),
        ("[>>]", "Умение писать промпты", "Общение с AI", False),
        ("[#]", "ГОТОВЫЙ ПРОЕКТ", "Калькулятор!", True),
    ]

    positions = [
        (CONTENT_LEFT, Inches(1.6)),
        (Inches(7), Inches(1.6)),
        (CONTENT_LEFT, Inches(3.8)),
        (Inches(7), Inches(3.8)),
    ]

    for i, ((icon, title, desc, highlighted), (left, top)) in enumerate(zip(results, positions)):
        if highlighted:
            card = add_card(
                slide, left, top,
                Inches(5.5), Inches(1.8),
                fill_color=ACID_15, border_color=CYBER_ACID
            )
        else:
            card = add_card(
                slide, left, top,
                Inches(5.5), Inches(1.8)
            )

        # Icon
        icon_color = CYBER_ACID if highlighted else TECH_GRAY
        add_text_box(
            slide, left + Inches(0.3), top + Inches(0.2),
            Inches(1), Inches(0.5), icon,
            font_name='Courier New', font_size=24,
            font_color=icon_color
        )

        # Title
        title_color = CYBER_ACID if highlighted else HOLO_WHITE
        add_text_box(
            slide, left + Inches(0.3), top + Inches(0.7),
            Inches(5), Inches(0.5), title,
            font_name='Arial', font_size=24,
            font_color=title_color, font_bold=True
        )

        # Description
        add_text_box(
            slide, left + Inches(0.3), top + Inches(1.2),
            Inches(5), Inches(0.4), desc,
            font_name='Arial', font_size=18,
            font_color=TECH_GRAY
        )

    # Footer
    add_footer(slide, "VIBE CODING STARTER")

    # Ready indicator
    add_text_box(
        slide, Inches(8.5), SLIDE_HEIGHT - Inches(0.6),
        Inches(4), Inches(0.4), "+ ГОТОВ К МОДУЛЮ 2",
        font_name='Courier New', font_size=16,
        font_color=CYBER_ACID, alignment=PP_ALIGN.RIGHT
    )

    return slide

def create_slide_11_lets_go(prs):
    """Slide 11: Let's Go!"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CYBER_VOID)

    # Subtitle
    add_text_box(
        slide, Inches(0), Inches(1.5),
        SLIDE_WIDTH, Inches(0.5), "МОДУЛЬ 1 • БЫСТРЫЙ СТАРТ",
        font_name='Courier New', font_size=24,
        font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
    )

    # Main title
    add_text_box(
        slide, Inches(0), Inches(2.3),
        SLIDE_WIDTH, Inches(1.5), "ПОЕХАЛИ!",
        font_name='Arial Black', font_size=120,
        font_color=HOLO_WHITE, alignment=PP_ALIGN.CENTER,
        font_bold=True
    )

    # Tagline
    add_text_box(
        slide, Inches(0), Inches(4),
        SLIDE_WIDTH, Inches(0.5), "Начинаем с урока 1.1",
        font_name='Arial', font_size=32,
        font_color=TECH_GRAY, alignment=PP_ALIGN.CENTER
    )

    # CTA Button
    button_width = Inches(5)
    button_left = (SLIDE_WIDTH - button_width) / 2

    cta_button = add_shape_with_text(
        slide, button_left, Inches(4.8),
        button_width, Inches(0.7), "ЧТО ТАКОЕ VIBE CODING ->",
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        fill_color=CYBER_ACID,
        font_name='Arial', font_size=22, font_bold=True,
        font_color=CYBER_VOID
    )

    # Progress dots (6 lessons)
    dots_text = "* o o o o o"  # First active, rest inactive
    add_text_box(
        slide, Inches(0), Inches(5.8),
        SLIDE_WIDTH, Inches(0.4), dots_text,
        font_name='Arial', font_size=24,
        font_color=CYBER_ACID, alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Inches(0), Inches(6.2),
        SLIDE_WIDTH, Inches(0.3), "(6 уроков)",
        font_name='Arial', font_size=16,
        font_color=TECH_GRAY, alignment=PP_ALIGN.CENTER
    )

    return slide

# ============================================================================
# MAIN
# ============================================================================

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create all slides
    print("Creating Slide 1: Title...")
    create_slide_1_title(prs)

    print("Creating Slide 2: Overview...")
    create_slide_2_overview(prs)

    print("Creating Slide 3: Program...")
    create_slide_3_program(prs)

    print("Creating Slide 4: Lesson 1.1...")
    create_slide_4_lesson_1_1(prs)

    print("Creating Slide 5: Lesson 1.2...")
    create_slide_5_lesson_1_2(prs)

    print("Creating Slide 6: Lesson 1.3...")
    create_slide_6_lesson_1_3(prs)

    print("Creating Slide 7: Lesson 1.4...")
    create_slide_7_lesson_1_4(prs)

    print("Creating Slide 8: Lesson 1.5...")
    create_slide_8_lesson_1_5(prs)

    print("Creating Slide 9: Lesson 1.6...")
    create_slide_9_lesson_1_6(prs)

    print("Creating Slide 10: Results...")
    create_slide_10_results(prs)

    print("Creating Slide 11: Let's Go...")
    create_slide_11_lets_go(prs)

    # Save
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "VIBE_CODING_MODULE1_INTRO.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")

    return output_path

if __name__ == "__main__":
    create_presentation()
